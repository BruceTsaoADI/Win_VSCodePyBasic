import os
from pptx import Presentation

# Constants: specify the root directory to search
INPUT_DIR = r"C:\Users\btsao\OneDrive - Analog Devices, Inc\Documents\BruceTsao\02_Document\ADI\WeeklyReport"
# To use a different folder, comment out above and uncomment below:

# Derive OUTPUT_FILE path to be in the same directory as this script
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_FILE = os.path.join(SCRIPT_DIR, "aggregated_ppt.txt")


def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_runs.append(shape.text)
    return "\n".join(text_runs)


def process_directory(input_dir, output_file):
    all_contents = []
    print(f"Searching in: {input_dir}")
    for dirpath, dirnames, filenames in os.walk(input_dir):
        for filename in filenames:
            if filename.lower().endswith(".pptx"):
                file_path = os.path.join(dirpath, filename)
                print(f"Processing file: {file_path}")
                try:
                    content = extract_text_from_pptx(file_path)
                    header = f"--- {file_path} ---"
                    all_contents.append(header)
                    all_contents.append(content)
                except Exception as e:
                    print(f"Error processing {file_path}: {e}")
    # Write all extracted text into a single file in the script directory
    try:
        with open(output_file, "w", encoding="utf-8") as f:
            f.write("\n\n".join(all_contents))
        print(f"All PPT contents saved to {output_file}")
    except PermissionError:
        print(
            f"Permission denied: cannot write to {output_file}. Please check the path and permissions."
        )


if __name__ == "__main__":
    process_directory(INPUT_DIR, OUTPUT_FILE)
    # read_ptt_files(r"C:\your\ptt\folder", "all_ptt_contents.txt")  # 使用範例
