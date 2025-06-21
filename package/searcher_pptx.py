import os
from pptx import Presentation
import csv


def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    slide_texts = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        slide_texts.append((i + 1, "\n".join(texts)))
    return slide_texts


def search_pptx_in_directory(
    directory, keywords, save_result=True, output_csv="pptx_keyword_results.csv"
):
    if isinstance(keywords, str):
        keywords = [keywords]  # allow single keyword as string

    matched_results = []
    for root, _, files in os.walk(directory):
        for filename in files:
            if filename.endswith(".pptx"):
                file_path = os.path.join(root, filename)
                try:
                    slides = extract_text_from_pptx(file_path)
                    for slide_no, content in slides:
                        for keyword in keywords:
                            if keyword.lower() in content.lower():
                                matched_results.append(
                                    [file_path, slide_no, keyword, content]
                                )
                                break  # avoid duplicate entries for multiple matching keywords
                except Exception as e:
                    print(f"Error reading {file_path}: {e}")

    if save_result:
        with open(output_csv, mode="w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(["File Path", "Slide No", "Matched Keyword", "Content"])
            writer.writerows(matched_results)
        print(
            f"✅ Done! Found {len(matched_results)} slides. Results saved to '{os.path.abspath(output_csv)}'."
        )
    else:
        print(f"✅ Done! Found {len(matched_results)} slides. No results saved.")


# 🔧 使用方式
# 修改以下兩行即可：
# 1️⃣ 指定你的 PPT 檔案資料夾
# 2️⃣ 指定要搜尋的關鍵字（可以是 list）


if __name__ == "__main__":
    DIRECTORY = r"C:\Users\btsao\OneDrive - Analog Devices, Inc\Documents\BruceTsao\02_Document\ADI\WeeklyReport"  # 替換成你的資料夾路徑
    KEYWORD = ["gmsl", "audio"]  # 替換成你要搜尋的關鍵字
    SAVE_RESULT = False
    search_pptx_in_directory(DIRECTORY, KEYWORD, SAVE_RESULT)
